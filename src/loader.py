import os
from pptx import Presentation
from .config import DATA_PATH


def load_slides():

    documents = []
    metadata = []

    os.makedirs("cache/extracted_slides", exist_ok=True)

    for week in os.listdir(DATA_PATH):

        week_path = os.path.join(DATA_PATH, week)

        if not os.path.isdir(week_path):
            continue

        week_txt_path = f"cache/extracted_slides/{week}.txt"

        with open(week_txt_path, "w", encoding="utf-8") as week_txt:

            for file in os.listdir(week_path):

                if not file.endswith(".pptx"):
                    continue

                file_path = os.path.join(week_path, file)

                prs = Presentation(file_path)

                for i, slide in enumerate(prs.slides):

                    text = []

                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)

                    slide_text = "\n".join(text).strip()

                    if len(slide_text) < 20:
                        continue

                    # save to txt
                    week_txt.write("\n\n------------------\n")
                    week_txt.write(f"{file} | slide {i+1}\n")
                    week_txt.write(slide_text)

                    documents.append(slide_text)

                    metadata.append({
                        "week": week,
                        "file": file,
                        "slide": i + 1
                    })

    return documents, metadata